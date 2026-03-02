from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Define Colors
    NXSTEP_BLUE = RGBColor(10, 25, 47)  # Slate-900 equivalent
    NXSTEP_ACCENT = RGBColor(14, 165, 233) # Sky-500 (Primary)
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(240, 240, 240)

    def set_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = NXSTEP_BLUE

    def add_title_slide(prs, title_text, subtitle_text):
        slide_layout = prs.slide_layouts[0] # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        set_background(slide)

        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = WHITE
        title.text_frame.paragraphs[0].font.bold = True
        
        subtitle = slide.placeholders[1]
        subtitle.text = subtitle_text
        subtitle.text_frame.paragraphs[0].font.color.rgb = LIGHT_GRAY

    def add_content_slide(prs, title_text, content_items):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        set_background(slide)

        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = NXSTEP_ACCENT
        title.text_frame.paragraphs[0].font.bold = True

        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.word_wrap = True

        for item in content_items:
            p = tf.add_paragraph()
            p.text = item
            p.font.color.rgb = WHITE
            p.font.size = Pt(20)
            p.space_after = Pt(14)

    # Slide 1: Title
    add_title_slide(prs, "NXSTEP", "Propelling Your Business Into the Digital Future")

    # Slide 2: The Problem
    add_content_slide(prs, "The Challenge: Digital Efficiency", [
        "Businesses struggle with fragmented workflows and manual data entry.",
        "Team productivity is lost on repetitive, low-value tasks.",
        "Disconnect between 300+ essential tools (CRM, Email, Slack).",
        "Scaling operations becomes impossible without automation."
    ])

    # Slide 3: The Solution
    add_content_slide(prs, "Our Solution: AI-Driven Automation", [
        "Seamless integration of your entire tech stack.",
        "Custom AI agents to handle complex workflows 24/7.",
        "Automated data syncing across all platforms.",
        "Reduction of manual errors by up to 99%."
    ])

    # Slide 4: Our Services
    add_content_slide(prs, "Services Tailored for Growth", [
        "Custom Automation Development (n8n, Make, Zapier).",
        "AI Agent Implementation & Training.",
        "Full-Stack Web Development (Next.js, React).",
        "Data Integration & Dashboarding."
    ])

    # Slide 5: Why NXSTEP?
    add_content_slide(prs, "Why Choose NXSTEP?", [
        "Proven Expertise: Specialists in both Web & AI.",
        "Efficiency First: We aim for 40% efficiency gains for clients.",
        "Rapid Deployment: Get up and running in weeks, not months.",
        "Dedicated Support: We partner with you for the long haul."
    ])

    # Slide 6: Case Study / Impact
    add_content_slide(prs, "Impact & Results", [
        "Client: E-commerce Retailer",
        "Challenge: Manual order processing took 4 hours/day.",
        "Solution: Automated order-to-fulfillment workflow.",
        "Result: 100% automation, 0 manual hours, 20% sales increase."
    ])

    # Slide 7: Next Steps
    add_title_slide(prs, "Ready to Transform?", "Contact us at contact@nxstep.fr\nVisit www.nxstep.fr")

    prs.save('NXSTEP_Business_Proposal.pptx')
    print("Presentation saved successfully.")

if __name__ == "__main__":
    create_presentation()
